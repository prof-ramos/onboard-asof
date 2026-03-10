from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x06, 0x5A, 0x82)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x66, 0x74)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
ACCENT = RGBColor(0x00, 0xA8, 0x96)
GOLD = RGBColor(0xF9, 0xE7, 0x95)
CORAL = RGBColor(0xF9, 0x61, 0x67)
ROOT_DIR = Path(__file__).resolve().parent.parent
OUTPUT_PATH = ROOT_DIR / "presentation" / "onboard-asof.pptx"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, *, size=20, color=INK,
             bold=False, font_name="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if text else [""]
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, fill, radius_text=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    if radius_text:
        shape.adjustments[0] = radius_text
    return shape


def add_circle_stat(slide, left, top, diameter, number, label, fill):
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, left, top, diameter, diameter
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.color.rgb = fill
    add_text(
        slide, number, left, top + Inches(0.28), diameter, Inches(0.45),
        size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Aptos Display"
    )
    add_text(
        slide, label, left - Inches(0.2), top + diameter + Inches(0.12),
        diameter + Inches(0.4), Inches(0.6), size=14, color=INK,
        align=PP_ALIGN.CENTER
    )


def add_title(slide, kicker, title, *, dark=False):
    kicker_color = ICE if dark else TEAL
    title_color = WHITE if dark else NAVY
    body_color = ICE if dark else MUTED
    add_text(slide, kicker.upper(), Inches(0.8), Inches(0.45), Inches(3.4), Inches(0.3),
             size=12, color=kicker_color, bold=True)
    add_text(slide, title, Inches(0.8), Inches(0.8), Inches(8.8), Inches(1.05),
             size=28, color=title_color, bold=True, font_name="Aptos Display")
    return body_color


def add_footer(slide, idx, total, dark=False):
    color = ICE if dark else MUTED
    add_text(slide, f"{idx:02d} / {total:02d}", Inches(11.82), Inches(6.72), Inches(0.78), Inches(0.25),
             size=10, color=color, align=PP_ALIGN.RIGHT)


def add_timeline_node(slide, left, top, label, sublabel, fill):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = fill
    circ.line.color.rgb = fill
    add_text(slide, label, left - Inches(0.2), top + Inches(0.55), Inches(1.2), Inches(0.25),
             size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sublabel, left - Inches(0.45), top + Inches(0.85), Inches(1.7), Inches(0.5),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_section_slide(slide, section, title, subtitle):
    set_bg(slide, NAVY)
    add_text(slide, section.upper(), Inches(0.9), Inches(1.0), Inches(3.5), Inches(0.35),
             size=14, color=GOLD, bold=True)
    add_text(slide, title, Inches(0.9), Inches(1.6), Inches(7.5), Inches(1.0),
             size=31, color=WHITE, bold=True, font_name="Aptos Display")
    add_text(slide, subtitle, Inches(0.9), Inches(3.0), Inches(6.2), Inches(0.8),
             size=18, color=ICE)
    add_card(slide, Inches(8.8), Inches(1.2), Inches(2.8), Inches(4.6), TEAL)
    add_text(slide, "Catálogo\ncompleto", Inches(9.15), Inches(2.0), Inches(2.1), Inches(0.9),
             size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Aptos Display")
    add_text(slide, "Listas por categoria para consulta rápida durante o onboard ou como anexo de apoio.", Inches(9.2), Inches(3.55),
             Inches(2.0), Inches(1.4), size=14, color=WHITE, align=PP_ALIGN.CENTER)


def add_list_panel(slide, title, items, left, top, width, height, *, fill=LIGHT, title_color=NAVY, text_color=INK):
    add_card(slide, left, top, width, height, fill)
    add_text(slide, title, left + Inches(0.22), top + Inches(0.22), width - Inches(0.44), Inches(0.28),
             size=16, color=title_color, bold=True)
    add_text(slide, "\n".join(items), left + Inches(0.22), top + Inches(0.62), width - Inches(0.44), height - Inches(0.84),
             size=13, color=text_color)


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 16

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_text(slide, "ONBOARD ASOF", Inches(0.8), Inches(0.7), Inches(4.5), Inches(0.5),
             size=16, color=ICE, bold=True)
    add_text(slide, "A associação criada pela carreira para representar, apoiar e conectar Oficiais de Chancelaria.",
             Inches(0.8), Inches(1.3), Inches(6.5), Inches(1.5),
             size=28, color=WHITE, bold=True, font_name="Aptos Display")
    add_text(slide, "Origem, missão, convênios, cadastro e contribuição associativa",
             Inches(0.8), Inches(3.0), Inches(5.0), Inches(0.55), size=17, color=WHITE, bold=True)
    add_text(slide, "1990  •  26 participantes  •  votação 24 x 2", Inches(0.82), Inches(4.22), Inches(5.6), Inches(0.35),
             size=15, color=ICE, bold=True)
    add_card(slide, Inches(8.95), Inches(1.05), Inches(3.15), Inches(4.78), TEAL)
    add_text(slide, "1990", Inches(9.28), Inches(1.45), Inches(1.2), Inches(0.5),
             size=28, color=WHITE, bold=True, font_name="Aptos Display")
    add_text(slide, "Fundação aprovada\npor 24 votos a 2", Inches(9.28), Inches(2.08),
             Inches(2.18), Inches(0.62), size=15, color=WHITE, bold=True)
    add_text(slide, "26 Oficiais de Chancelaria, reunidos em Brasília, decidiram criar uma entidade própria e independente.",
             Inches(9.28), Inches(3.02), Inches(2.18), Inches(1.0), size=13, color=WHITE)
    add_text(slide, "Mensagem central", Inches(9.28), Inches(4.35), Inches(1.5), Inches(0.22),
             size=11, color=GOLD, bold=True)
    add_text(slide, "Representação institucional com apoio prático ao associado.", Inches(9.28), Inches(4.62),
             Inches(2.18), Inches(0.55), size=13, color=WHITE, bold=True)
    add_footer(slide, 1, total, dark=True)

    # 2. Origin timeline
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Origem", "A ASOF nasce de uma decisão coletiva da carreira", dark=False)
    add_text(slide, "Em 7 de agosto de 1990, Oficiais de Chancelaria decidiram criar uma associação própria para representar a carreira com autonomia institucional.",
             Inches(0.8), Inches(1.75), Inches(7.0), Inches(0.8), size=18, color=body_color)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(4.0), Inches(11.6), Inches(4.0))
    line.line.color.rgb = ICE
    line.line.width = Pt(3)
    add_timeline_node(slide, Inches(1.35), Inches(3.78), "7 ago 1990", "Reunião inicial", TEAL)
    add_timeline_node(slide, Inches(4.35), Inches(3.78), "26", "participantes", ACCENT)
    add_timeline_node(slide, Inches(7.3), Inches(3.78), "24 x 2", "votação", CORAL)
    add_timeline_node(slide, Inches(10.2), Inches(3.78), "ASOF", "fundação imediata", NAVY)
    add_card(slide, Inches(8.8), Inches(1.45), Inches(3.6), Inches(1.5), WHITE)
    add_text(slide, "Brasília, ASMRE, 18h", Inches(9.1), Inches(1.75), Inches(2.8), Inches(0.25),
             size=17, color=NAVY, bold=True)
    add_text(slide, "O debate considerou manter apenas um núcleo em outra entidade, mas prevaleceu a necessidade de representação própria.",
             Inches(9.1), Inches(2.1), Inches(2.8), Inches(0.65), size=13, color=MUTED)
    add_footer(slide, 2, total)

    # 3. Autonomy choice
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    body_color = add_title(slide, "Escolha central", "A autonomia institucional foi a opção que definiu a associação", dark=False)
    add_text(slide, "A fundação da ASOF responde a uma necessidade específica: representar os interesses da carreira com identidade própria e personalidade jurídica independente.",
             Inches(0.8), Inches(1.8), Inches(7.2), Inches(0.8), size=18, color=body_color)
    add_card(slide, Inches(0.9), Inches(3.0), Inches(5.55), Inches(2.7), LIGHT)
    add_text(slide, "Alternativa considerada", Inches(1.2), Inches(3.3), Inches(2.8), Inches(0.3),
             size=14, color=TEAL, bold=True)
    add_text(slide, "Manter apenas um núcleo de Oficiais de Chancelaria dentro de uma estrutura associativa já existente.", Inches(1.2), Inches(3.75),
             Inches(4.7), Inches(1.1), size=20, color=INK, bold=True)
    add_text(slide, "Preservaria apoio inicial, mas não resolveria a necessidade de identidade institucional própria.", Inches(1.2), Inches(5.0),
             Inches(4.7), Inches(0.45), size=13, color=MUTED)
    add_card(slide, Inches(6.85), Inches(3.0), Inches(5.55), Inches(2.7), NAVY)
    add_text(slide, "Opção aprovada", Inches(7.15), Inches(3.3), Inches(2.0), Inches(0.3),
             size=14, color=GOLD, bold=True)
    add_text(slide, "Criar uma associação exclusiva da categoria, com foco integral nos interesses dos Oficiais de Chancelaria.", Inches(7.15), Inches(3.75),
             Inches(4.75), Inches(1.15), size=20, color=WHITE, bold=True)
    add_text(slide, "A decisão consolidou autonomia, representação específica e continuidade institucional.", Inches(7.15), Inches(5.0),
             Inches(4.75), Inches(0.45), size=13, color=ICE)
    add_footer(slide, 3, total)

    # 4. Mission pillars
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Missão", "O estatuto transforma a fundação em compromisso permanente", dark=False)
    add_text(slide, "O estatuto consolida a ASOF como entidade civil sem fins lucrativos voltada à união da carreira, à representação institucional e à geração de apoio prático ao associado.",
             Inches(0.8), Inches(1.75), Inches(8.0), Inches(0.75), size=18, color=body_color)
    cards = [
        (Inches(0.95), Inches(3.0), TEAL, "União da carreira"),
        (Inches(3.95), Inches(3.0), NAVY, "Representação institucional"),
        (Inches(6.95), Inches(3.0), ACCENT, "Aprimoramento profissional"),
        (Inches(9.95), Inches(3.0), CORAL, "Benefícios e convênios"),
    ]
    for left, top, color, text in cards:
        add_card(slide, left, top, Inches(2.3), Inches(2.25), color)
        add_text(slide, text, left + Inches(0.25), top + Inches(0.75), Inches(1.8), Inches(0.8),
                 size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4, total)

    # 5. Value system
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    body_color = add_title(slide, "Entrega de valor", "O valor da ASOF combina representação, organização e apoio ao associado", dark=False)
    add_text(slide, "A atuação associativa não se limita à defesa institucional. Ela também organiza o vínculo do associado e viabiliza benefícios concretos para a carreira e seus dependentes.",
             Inches(0.8), Inches(1.8), Inches(7.3), Inches(0.7), size=18, color=body_color)
    steps = [
        ("1", "Representar", "Leva demandas da carreira a autoridades e instituições."),
        ("2", "Organizar", "Mantém base cadastral, regras e relacionamento associativo."),
        ("3", "Conectar", "Articula convênios e benefícios em diferentes frentes."),
        ("4", "Apoiar", "Transforma a associação em suporte contínuo ao associado."),
    ]
    lefts = [Inches(0.95), Inches(3.05), Inches(5.15), Inches(7.25)]
    colors = [TEAL, NAVY, ACCENT, CORAL]
    for i, (num, title, desc) in enumerate(steps):
        add_card(slide, lefts[i], Inches(3.1), Inches(1.8), Inches(2.7), LIGHT)
        add_circle_stat(slide, lefts[i] + Inches(0.57), Inches(3.3), Inches(0.65), num, "", colors[i])
        add_text(slide, title, lefts[i] + Inches(0.18), Inches(4.15), Inches(1.45), Inches(0.3),
                 size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, lefts[i] + Inches(0.15), Inches(4.55), Inches(1.5), Inches(0.9),
                 size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 5, total)

    # 6. Convenios
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Benefícios", "Convênios ampliam o valor da associação no dia a dia", dark=True)
    add_text(slide, "O acervo disponível mostra uma rede ampla de parcerias que leva a atuação da ASOF para educação, saúde, bem-estar e serviços.", Inches(0.8), Inches(1.75),
             Inches(7.5), Inches(0.7), size=18, color=ICE)
    add_circle_stat(slide, Inches(1.0), Inches(3.0), Inches(1.15), "28", "Educação", TEAL)
    add_circle_stat(slide, Inches(3.25), Inches(3.0), Inches(1.15), "24", "Saúde", ACCENT)
    add_circle_stat(slide, Inches(5.5), Inches(3.0), Inches(1.15), "25", "Lazer e bem-estar", CORAL)
    add_circle_stat(slide, Inches(7.75), Inches(3.0), Inches(1.15), "5", "Serviços", GOLD)
    add_card(slide, Inches(10.0), Inches(2.55), Inches(2.1), Inches(2.65), WHITE)
    add_text(slide, "Mais que desconto", Inches(10.25), Inches(2.95), Inches(1.6), Inches(0.3),
             size=14, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Os convênios ajudam a transformar presença associativa em valor percebido no cotidiano.", Inches(10.25), Inches(3.4),
             Inches(1.6), Inches(1.1), size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6, total, dark=True)

    # 7. Convenio flow
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Funcionamento", "Os convênios operam como uma rede de acesso organizada pela ASOF", dark=False)
    add_text(slide, "A associação articula o benefício, identifica o associado e divulga as parcerias. A prestação do serviço permanece com a empresa conveniada.", Inches(0.8), Inches(1.8),
             Inches(7.9), Inches(0.65), size=18, color=body_color)
    for left, label, fill in [
        (Inches(1.0), "ASOF", TEAL),
        (Inches(4.7), "Associado\ne dependentes", NAVY),
        (Inches(8.65), "Empresa\nconveniada", ACCENT),
    ]:
        add_card(slide, left, Inches(3.15), Inches(2.7), Inches(1.75), fill)
        add_text(slide, label, left + Inches(0.25), Inches(3.65), Inches(2.2), Inches(0.6),
                 size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for start, end in [(Inches(3.72), Inches(4.7)), (Inches(7.4), Inches(8.65))]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start, Inches(4.0), end, Inches(4.0))
        conn.line.color.rgb = ICE
        conn.line.width = Pt(2.5)
    add_text(slide, "declaração ou carteirinha", Inches(3.2), Inches(3.5), Inches(1.6), Inches(0.3),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "prestação do serviço", Inches(7.45), Inches(3.5), Inches(1.5), Inches(0.3),
             size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 7, total)

    # 8. Form
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    body_color = add_title(slide, "Cadastro", "O vínculo associativo depende de dados claros e atualizados", dark=False)
    add_text(slide, "O formulário de atualização de dados organiza informações pessoais, funcionais, bancárias e de contato, sustentando comunicação, atendimento e cobrança associativa.", Inches(0.8), Inches(1.8),
             Inches(8.1), Inches(0.7), size=18, color=body_color)
    steps = [("1", "Preencher"), ("2", "Assinar"), ("3", "Remeter à ASOF")]
    step_lefts = [Inches(1.05), Inches(4.25), Inches(7.45)]
    for i, (num, label) in enumerate(steps):
        add_card(slide, step_lefts[i], Inches(3.1), Inches(2.35), Inches(2.1), LIGHT)
        add_circle_stat(slide, step_lefts[i] + Inches(0.86), Inches(3.35), Inches(0.6), num, "", [TEAL, NAVY, ACCENT][i])
        add_text(slide, label, step_lefts[i] + Inches(0.22), Inches(4.32), Inches(1.9), Inches(0.35),
                 size=19, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(10.3), Inches(2.9), Inches(2.0), Inches(2.5), TEAL)
    add_text(slide, "Dados-chave", Inches(10.55), Inches(3.2), Inches(1.5), Inches(0.25),
             size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "contatos\nsituação funcional\ndependentes\ndados bancários", Inches(10.55), Inches(3.65),
             Inches(1.5), Inches(1.2), size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8, total)

    # 9. Contribution
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Sustentação", "A contribuição associativa financia a continuidade institucional da ASOF", dark=False)
    add_text(slide, "O estatuto prevê a contribuição mensal como fonte regular de receita. O formulário atual mostra os valores praticados para desconto em folha.", Inches(0.8), Inches(1.8),
             Inches(7.7), Inches(0.65), size=18, color=body_color)
    table = [
        (Inches(1.1), TEAL, "Ativo", "R$ 40,00"),
        (Inches(4.2), NAVY, "Aposentado", "R$ 20,00"),
        (Inches(7.3), ACCENT, "Exterior", "US$ 25,00"),
    ]
    for left, fill, label, value in table:
        add_card(slide, left, Inches(3.0), Inches(2.45), Inches(2.2), fill)
        add_text(slide, label, left + Inches(0.2), Inches(3.4), Inches(2.05), Inches(0.35),
                 size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, value, left + Inches(0.2), Inches(4.0), Inches(2.05), Inches(0.5),
                 size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Aptos Display")
    add_text(slide, "Leitura útil para o onboard:\n\nEstatuto: define a contribuição como base regular de receita.\nFormulário: explicita os valores operacionais do documento em uso.", Inches(10.2), Inches(2.85),
             Inches(2.0), Inches(2.0), size=14, color=MUTED)
    add_footer(slide, 9, total)

    # 10. Closing
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Fechamento", "Associar-se é entrar em uma estrutura de representação e suporte", dark=True)
    add_text(slide, "A ASOF conecta identidade profissional, ação coletiva e benefícios concretos em uma mesma estrutura associativa.", Inches(0.8), Inches(1.8),
             Inches(7.4), Inches(0.7), size=20, color=ICE)
    pillars = [
        (Inches(0.95), TEAL, "Representar\na carreira"),
        (Inches(4.35), ACCENT, "Fortalecer\nvínculos"),
        (Inches(7.75), CORAL, "Gerar apoio\nprático"),
    ]
    for left, fill, text in pillars:
        add_card(slide, left, Inches(3.05), Inches(2.7), Inches(2.0), fill)
        add_text(slide, text, left + Inches(0.2), Inches(3.65), Inches(2.3), Inches(0.7),
                 size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(10.75), Inches(2.85), Inches(1.25), Inches(2.35), WHITE)
    add_text(slide, "Desde\n1990", Inches(10.9), Inches(3.4), Inches(0.95), Inches(0.7),
             size=25, color=NAVY, bold=True, align=PP_ALIGN.CENTER, font_name="Aptos Display")
    add_text(slide, "Criada pela carreira\npara servir à carreira", Inches(10.75), Inches(4.45),
             Inches(1.25), Inches(0.6), size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 10, total, dark=True)

    # 11. Appendix divider
    slide = prs.slides.add_slide(blank)
    add_section_slide(
        slide,
        "Apêndice",
        "Lista completa dos convênios",
        "Os slides a seguir organizam os parceiros por categoria, com foco em consulta rápida e melhor legibilidade."
    )
    add_footer(slide, 11, total, dark=True)

    # 12. Convenios - Educacao
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    body_color = add_title(slide, "Apêndice", "Convênios de educação", dark=False)
    add_text(slide, "Parceiros identificados no acervo consultado para a frente de educação.", Inches(0.8), Inches(1.75),
             Inches(7.0), Inches(0.45), size=17, color=body_color)
    add_list_panel(slide, "Educação 1", [
        "Aliança Francesa", "CCAA", "CIMAN", "Colégio Dinatos COC", "Cultura Inglesa",
        "Curso Ignis", "Escola Franciscana Fátima", "Escola Presbiteriana Mackenzie", "FACITEC"
    ], Inches(0.85), Inches(2.45), Inches(3.8), Inches(3.9))
    add_list_panel(slide, "Educação 2", [
        "FGV", "Gran Cursos", "IBMEC", "IDP", "IESB",
        "Instituto Blaise Pascal", "Laboro", "MS Educação", "Positive Idiomas"
    ], Inches(4.85), Inches(2.45), Inches(3.8), Inches(3.9))
    add_list_panel(slide, "Educação 3", [
        "Rede de Ensino JK", "St. Giles", "Studio On-line", "Swiss International School",
        "Thomas Jefferson", "UDF", "UniCEUB", "UNIP", "Wizard"
    ], Inches(8.85), Inches(2.45), Inches(3.6), Inches(3.9))
    add_list_panel(slide, "Complementares", [
        "Centro Educacional Maria Auxiliadora", "UNYLEYA",
        "Instituto Formação para a Educação", "3W Educacional Editora e Cursos"
    ], Inches(8.85), Inches(6.55), Inches(3.6), Inches(0.55), fill=TEAL, title_color=WHITE, text_color=WHITE)
    add_footer(slide, 12, total)

    # 13. Convenios - Saude parte 1
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Apêndice", "Convênios de saúde", dark=False)
    add_text(slide, "Primeira parte da frente de saúde, separada para manter leitura confortável em tela.", Inches(0.8), Inches(1.75),
             Inches(8.2), Inches(0.45), size=17, color=body_color)
    add_list_panel(slide, "Saúde 1", [
        "Aquafisio Hidroterapia Esportiva", "Bonvena Medicina Reprodutiva", "CBV",
        "CDI Centro Diagnóstico por Imagem", "Centro de Medicina Nuclear da Guanabara",
        "CIORB", "Clínica Plenus Psicologia e Saúde Integrada", "Dermatologic",
        "EBM-Odonto", "Endogastrus", "Home Hospital Ortopédico e Medicina Especializada"
    ], Inches(0.85), Inches(2.45), Inches(5.5), Inches(4.35))
    add_list_panel(slide, "Saúde 2", [
        "Implantocard", "Juliana Cristina Paim Psicóloga Clínica e Neuropsicóloga",
        "Laboratório Citoprev", "Mei Li Acupuntura", "Oculare Oftalmologia",
        "Odontobrasília", "Odontoempresa", "Ressonance",
        "Rita Trindade", "RM Clínica Médica e Estética", "Sabin", "Souli Psicologia e Psicopedagogia"
    ], Inches(6.65), Inches(2.45), Inches(5.8), Inches(4.35))
    add_footer(slide, 13, total)

    # 14. Convenios - Saude complementar
    slide = prs.slides.add_slide(blank)
    set_bg(slide, WHITE)
    body_color = add_title(slide, "Apêndice", "Convênios de saúde complementar", dark=False)
    add_text(slide, "Parceiros adicionais citados em subpastas e documentação complementar do acervo.", Inches(0.8), Inches(1.75),
             Inches(7.6), Inches(0.45), size=17, color=body_color)
    add_list_panel(slide, "Complementares", [
        "Microsom", "Oftalmocenter", "Implantomed",
        "Instituto Oftalmológico Visão", "Olhar Hospital Oftalmológico Ltda."
    ], Inches(0.9), Inches(2.55), Inches(4.0), Inches(2.8), fill=NAVY, title_color=WHITE, text_color=WHITE)
    add_card(slide, Inches(5.35), Inches(2.55), Inches(6.95), Inches(2.8), TEAL)
    add_text(slide, "Leitura editorial", Inches(5.65), Inches(2.85), Inches(2.5), Inches(0.25),
             size=16, color=WHITE, bold=True)
    add_text(slide, "A frente de saúde é a mais densa do acervo e cobre diagnósticos, clínicas, odontologia, psicologia, oftalmologia e serviços terapêuticos.",
             Inches(5.65), Inches(3.25), Inches(6.2), Inches(1.35), size=19, color=WHITE, bold=True)
    add_text(slide, "Na apresentação ao vivo, vale citar poucos exemplos e usar este slide como apoio de consulta, não como leitura integral.", Inches(5.65), Inches(4.7),
             Inches(5.9), Inches(0.45), size=13, color=ICE)
    add_footer(slide, 14, total)

    # 15. Convenios - Lazer
    slide = prs.slides.add_slide(blank)
    set_bg(slide, LIGHT)
    body_color = add_title(slide, "Apêndice", "Convênios de lazer e bem-estar", dark=False)
    add_text(slide, "A categoria reúne atividades físicas, turismo, estética, hotelaria e serviços voltados à qualidade de vida.", Inches(0.8), Inches(1.75),
             Inches(8.3), Inches(0.45), size=17, color=body_color)
    add_list_panel(slide, "Lazer e bem-estar 1", [
        "Academia Club 22", "Academia Companhia Athletica", "Academia Dom Bosco",
        "Academia Runway", "Academia Team Nogueira", "Academia Vasco Neto",
        "ASBAC", "Associação Cristã de Moços", "ATP Atividades Físicas",
        "Bancorbrás Consórcios", "Bancorbrás Turismo", "Bay Park Resort Hotel"
    ], Inches(0.85), Inches(2.45), Inches(5.6), Inches(4.35))
    add_list_panel(slide, "Lazer e bem-estar 2", [
        "Camisaria Nyll", "Clube dos Previdenciários", "Corpus Studio de Pilates",
        "Elia Spa", "Equilibriom Estética e Podologia", "Escola Ballet Garden",
        "GrandBittar Hotel", "Hotel Fazenda Mestre Darmas", "Montreal Turismo",
        "Nuwa Spa", "Studio Caracóis", "Studio Hair 180 Graus", "Thermas do Rio Quente"
    ], Inches(6.65), Inches(2.45), Inches(5.8), Inches(4.35))
    add_footer(slide, 15, total)

    # 16. Convenios - Servicos
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_title(slide, "Apêndice", "Convênios de serviços", dark=True)
    add_text(slide, "Parcerias ligadas ao apoio ao dia a dia do associado identificadas no acervo disponível.", Inches(0.8), Inches(1.75),
             Inches(7.1), Inches(0.5), size=17, color=ICE)
    service_cards = [
        (Inches(1.0), Inches(3.0), TEAL, "Hertz"),
        (Inches(3.4), Inches(3.0), ACCENT, "Lavanderia\nCleaners Club"),
        (Inches(5.8), Inches(3.0), CORAL, "Lavanderia\nSelecta BonaSecco"),
        (Inches(8.2), Inches(3.0), GOLD, "Óptica Elen"),
        (Inches(10.6), Inches(3.0), WHITE, "Grupo Porcão"),
    ]
    for left, top, fill, label in service_cards:
        add_card(slide, left, top, Inches(1.75), Inches(1.9), fill)
        text_color = NAVY if fill == WHITE or fill == GOLD else WHITE
        add_text(slide, label, left + Inches(0.15), top + Inches(0.58), Inches(1.45), Inches(0.7),
                 size=18, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Esses convênios complementam a proposta de valor da ASOF ao conectar representação institucional com utilidade cotidiana.", Inches(1.0), Inches(5.55),
             Inches(11.0), Inches(0.6), size=18, color=ICE, align=PP_ALIGN.CENTER)
    add_footer(slide, 16, total, dark=True)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))


if __name__ == "__main__":
    build_deck()
